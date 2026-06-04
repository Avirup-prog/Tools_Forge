"""
routers/pdf.py — 14 PDF tool endpoints
POST /api/pdf/to-word
POST /api/pdf/from-word
POST /api/pdf/to-excel
POST /api/pdf/to-ppt
POST /api/pdf/merge
POST /api/pdf/split
POST /api/pdf/compress
POST /api/pdf/encrypt
POST /api/pdf/decrypt
POST /api/pdf/rotate
POST /api/pdf/to-image
POST /api/pdf/from-image
POST /api/pdf/watermark
POST /api/pdf/fill-sign
"""

import io
import os
import json
import zipfile
import tempfile
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import Response, StreamingResponse

import pdfplumber
import pypdf
from pypdf import PdfReader, PdfWriter
from pypdf.generic import NameObject
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.colors import Color, black
from docx import Document
from PIL import Image as PILImage
import img2pdf
from pdf2docx import Converter
from fastapi.responses import FileResponse

from utils.helpers import save_upload, temp_path, cleanup, tmp_scope, require_mime, tool_error

router = APIRouter()


# ─────────────────────────────────────────────────────────────────────────────
# 1. PDF → Word  /api/pdf/to-word
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/to-word", summary="Convert PDF to DOCX")
async def pdf_to_word(file: UploadFile = File(...)):
    require_mime(file, "pdf")

    src = await save_upload(file, ".pdf")
    out = temp_path(".docx")

    try:
        cv = Converter(str(src))
        cv.convert(str(out))
        cv.close()
)

        data = out.read_bytes()

        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition":
                f'attachment; filename="{Path(file.filename).stem}.docx"'
            },
        )

    except Exception as e:
        raise tool_error(f"PDF to Word failed: {e}")

    finally:
        cleanup(src, out)

# ─────────────────────────────────────────────────────────────────────────────
# 2. Word → PDF  /api/pdf/from-word
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/from-word", summary="Convert DOCX to PDF")
async def word_to_pdf(file: UploadFile = File(...)):
    require_mime(file, "word")
    src = await save_upload(file, ".docx")
    out = temp_path(".pdf")

    try:
        # Read docx and produce PDF via reportlab
        doc = Document(str(src))
        buf = io.BytesIO()
        c = rl_canvas.Canvas(buf, pagesize=letter)
        width, height = letter
        margin = inch
        y = height - margin
        line_h = 14

        def new_page():
            nonlocal y
            c.showPage()
            y = height - margin

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                y -= line_h / 2
                if y < margin:
                    new_page()
                continue

            # Heading styles
            if para.style.name.startswith("Heading 1"):
                c.setFont("Helvetica-Bold", 16)
            elif para.style.name.startswith("Heading 2"):
                c.setFont("Helvetica-Bold", 13)
            elif para.style.name.startswith("Heading"):
                c.setFont("Helvetica-Bold", 11)
            else:
                c.setFont("Helvetica", 10)

            # Word-wrap
            words = text.split()
            line = ""
            max_w = width - 2 * margin
            for word in words:
                test = f"{line} {word}".strip()
                if c.stringWidth(test, c._fontname, c._fontsize) < max_w:
                    line = test
                else:
                    c.drawString(margin, y, line)
                    y -= line_h
                    if y < margin:
                        new_page()
                    line = word
            if line:
                c.drawString(margin, y, line)
                y -= line_h * 1.4
            if y < margin:
                new_page()

        c.save()
        buf.seek(0)
        data = buf.getvalue()
        return Response(
            content=data,
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Word to PDF failed: {e}")
    finally:
        cleanup(src, out)


# ─────────────────────────────────────────────────────────────────────────────
# 3. PDF → Excel  /api/pdf/to-excel
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/to-excel", summary="Extract PDF tables to XLSX")
async def pdf_to_excel(file: UploadFile = File(...)):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")
    out = temp_path(".xlsx")

    try:
        # Use openpyxl directly to avoid pandas dependency
        import openpyxl
        wb = openpyxl.Workbook()
        wb.remove(wb.active)   # remove default sheet

        with pdfplumber.open(str(src)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()
                if tables:
                    for t_idx, table in enumerate(tables, 1):
                        ws = wb.create_sheet(f"P{page_num}_T{t_idx}")
                        for row in table:
                            ws.append([str(cell) if cell is not None else "" for cell in row])
                else:
                    # Fallback: dump text lines as single-column sheet
                    ws = wb.create_sheet(f"Page {page_num}")
                    text = page.extract_text() or ""
                    for line in text.split("\n"):
                        if line.strip():
                            ws.append([line.strip()])

        if not wb.sheetnames:
            wb.create_sheet("Empty")

        wb.save(str(out))
        data = out.read_bytes()
        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}.xlsx"'},
        )
    except Exception as e:
        raise tool_error(f"PDF to Excel failed: {e}")
    finally:
        cleanup(src, out)


# ─────────────────────────────────────────────────────────────────────────────
# 4. PDF → PowerPoint  /api/pdf/to-ppt
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/to-ppt", summary="Convert PDF pages to PPTX slides")
async def pdf_to_ppt(file: UploadFile = File(...)):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")
    out = temp_path(".pptx")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]   # blank

        with pdfplumber.open(str(src)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                slide = prs.slides.add_slide(blank_layout)
                text = page.extract_text() or f"(Page {i} — no extractable text)"

                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
                tf = txBox.text_frame
                tf.word_wrap = True

                lines = [l.strip() for l in text.split("\n") if l.strip()]
                for j, line in enumerate(lines):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = line
                    run.font.size = Pt(11 if j > 0 else 14)
                    if j == 0:
                        run.font.bold = True

        prs.save(str(out))
        data = out.read_bytes()
        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}.pptx"'},
        )
    except Exception as e:
        raise tool_error(f"PDF to PPT failed: {e}")
    finally:
        cleanup(src, out)


# ─────────────────────────────────────────────────────────────────────────────
# 5. Merge PDFs  /api/pdf/merge
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/merge", summary="Merge multiple PDFs into one")
async def merge_pdf(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise tool_error("Please upload at least 2 PDF files to merge.")
    for f in files:
        require_mime(f, "pdf")

    srcs = []
    try:
        writer = PdfWriter()
        for f in files:
            path = await save_upload(f, ".pdf")
            srcs.append(path)
            reader = PdfReader(str(path))
            for page in reader.pages:
                writer.add_page(page)

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": 'attachment; filename="merged.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Merge failed: {e}")
    finally:
        cleanup(*srcs)


# ─────────────────────────────────────────────────────────────────────────────
# 6. Split PDF  /api/pdf/split
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/split", summary="Split PDF into individual pages or ranges")
async def split_pdf(
    file: UploadFile = File(...),
    ranges: Optional[str] = Form(None),   # e.g. "1-3,5,7-9"  (1-indexed)
):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")

    try:
        reader = PdfReader(str(src))
        total = len(reader.pages)

        # Parse ranges
        def parse_ranges(spec: str, total: int):
            pages = set()
            for part in spec.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    pages.update(range(int(a) - 1, min(int(b), total)))
                else:
                    n = int(part) - 1
                    if 0 <= n < total:
                        pages.add(n)
            return sorted(pages)

        if ranges:
            page_indices = parse_ranges(ranges, total)
        else:
            page_indices = list(range(total))

        # Build ZIP of individual PDFs
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for idx in page_indices:
                writer = PdfWriter()
                writer.add_page(reader.pages[idx])
                page_buf = io.BytesIO()
                writer.write(page_buf)
                zf.writestr(f"page_{idx + 1:03d}.pdf", page_buf.getvalue())

        return Response(
            content=zip_buf.getvalue(),
            media_type="application/zip",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_split.zip"'},
        )
    except ValueError as e:
        raise tool_error(f"Invalid range format: {e}")
    except Exception as e:
        raise tool_error(f"Split failed: {e}")
    finally:
        cleanup(src)


# ─────────────────────────────────────────────────────────────────────────────
# 7. Compress PDF  /api/pdf/compress
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/compress", summary="Reduce PDF file size")
async def compress_pdf(
    file: UploadFile = File(...),
    quality: int = Form(75),   # image quality 1-95
):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")

    try:
        reader = PdfReader(str(src))
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Compress embedded images in each page
        for page in writer.pages:
            for img_name in list(page.get("/Resources", {}).get("/XObject", {}).keys()):
                xobj = page["/Resources"]["/XObject"][img_name]
                if xobj.get("/Subtype") == "/Image":
                    try:
                        xobj.compress_content_streams()
                    except Exception:
                        pass

        # General stream compression
        for page in writer.pages:
            page.compress_content_streams()

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_compressed.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Compress failed: {e}")
    finally:
        cleanup(src)


# ─────────────────────────────────────────────────────────────────────────────
# 8. Encrypt (Add Password)  /api/pdf/encrypt
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/encrypt", summary="Password-protect a PDF")
async def add_password(
    file: UploadFile = File(...),
    password: str = Form(...),
    owner_password: Optional[str] = Form(None),
):
    require_mime(file, "pdf")
    if not password:
        raise tool_error("Password is required.")
    src = await save_upload(file, ".pdf")

    try:
        reader = PdfReader(str(src))
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(
            user_password=password,
            owner_password=owner_password or password,
            use_128bit=True,
        )

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_protected.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Encryption failed: {e}")
    finally:
        cleanup(src)


# ─────────────────────────────────────────────────────────────────────────────
# 9. Decrypt (Remove Password)  /api/pdf/decrypt
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/decrypt", summary="Remove password from a PDF")
async def remove_password(
    file: UploadFile = File(...),
    password: str = Form(...),
):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")

    try:
        reader = PdfReader(str(src))
        if reader.is_encrypted:
            result = reader.decrypt(password)
            if result == 0:
                raise tool_error("Wrong password — could not decrypt PDF.", 401)

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_unlocked.pdf"'},
        )
    except HTTPException:
        raise
    except Exception as e:
        raise tool_error(f"Decrypt failed: {e}")
    finally:
        cleanup(src)


# ─────────────────────────────────────────────────────────────────────────────
# 10. Rotate & Reorder  /api/pdf/rotate
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/rotate", summary="Rotate PDF pages")
async def rotate_pdf(
    file: UploadFile = File(...),
    angle: int = Form(90),           # 90, 180, 270
    pages: Optional[str] = Form(None),  # "all" or "1,3,5" (1-indexed)
):
    require_mime(file, "pdf")
    if angle not in (90, 180, 270):
        raise tool_error("Angle must be 90, 180, or 270.")
    src = await save_upload(file, ".pdf")

    try:
        reader = PdfReader(str(src))
        writer = PdfWriter()
        total = len(reader.pages)

        # Parse which pages to rotate
        if pages and pages.lower() != "all":
            rotate_set = set()
            for part in pages.split(","):
                part = part.strip()
                n = int(part) - 1
                if 0 <= n < total:
                    rotate_set.add(n)
        else:
            rotate_set = set(range(total))

        for i, page in enumerate(reader.pages):
            if i in rotate_set:
                page.rotate(angle)
            writer.add_page(page)

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_rotated.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Rotate failed: {e}")
    finally:
        cleanup(src)


# ─────────────────────────────────────────────────────────────────────────────
# 11. PDF → Image  /api/pdf/to-image
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/to-image", summary="Convert PDF pages to images (ZIP of PNGs)")
async def pdf_to_image(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    fmt: str = Form("png"),   # png or jpg
):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")

    try:
        from pdf2image import convert_from_path
        fmt = fmt.lower()
        if fmt not in ("png", "jpg", "jpeg"):
            fmt = "png"

        images = convert_from_path(str(src), dpi=dpi)

        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, img in enumerate(images, 1):
                img_buf = io.BytesIO()
                save_fmt = "JPEG" if fmt in ("jpg", "jpeg") else "PNG"
                img.save(img_buf, format=save_fmt, optimize=True)
                ext = "jpg" if fmt in ("jpg", "jpeg") else "png"
                zf.writestr(f"page_{i:03d}.{ext}", img_buf.getvalue())

        return Response(
            content=zip_buf.getvalue(),
            media_type="application/zip",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_images.zip"'},
        )
    except Exception as e:
        raise tool_error(f"PDF to image failed: {e}")
    finally:
        cleanup(src)


# ─────────────────────────────────────────────────────────────────────────────
# 12. Image → PDF  /api/pdf/from-image
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/from-image", summary="Convert one or more images to a PDF")
async def image_to_pdf(files: List[UploadFile] = File(...)):
    for f in files:
        require_mime(f, "image")

    srcs = []
    try:
        for f in files:
            p = await save_upload(f, Path(f.filename or "img").suffix or ".jpg")
            # Ensure image is RGB (no alpha for PDF/JPEG compat)
            img = PILImage.open(str(p)).convert("RGB")
            img.save(str(p), format="JPEG")
            srcs.append(p)

        pdf_bytes = img2pdf.convert([str(p) for p in srcs])
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": 'attachment; filename="images.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Image to PDF failed: {e}")
    finally:
        cleanup(*srcs)


# ─────────────────────────────────────────────────────────────────────────────
# 13. Add Watermark  /api/pdf/watermark
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/watermark", summary="Add text watermark to every PDF page")
async def add_watermark(
    file: UploadFile = File(...),
    text: str = Form("CONFIDENTIAL"),
    opacity: float = Form(0.15),
    font_size: int = Form(48),
    color: str = Form("gray"),   # gray | red | blue
):
    require_mime(file, "pdf")
    src = await save_upload(file, ".pdf")
    wm_path = temp_path(".pdf")

    try:
        reader = PdfReader(str(src))
        first_page = reader.pages[0]
        pw = float(first_page.mediabox.width)
        ph = float(first_page.mediabox.height)

        # Build watermark PDF with reportlab
        wm_buf = io.BytesIO()
        c = rl_canvas.Canvas(wm_buf, pagesize=(pw, ph))

        color_map = {"red": Color(1, 0, 0, alpha=opacity),
                     "blue": Color(0, 0, 1, alpha=opacity),
                     "gray": Color(0.5, 0.5, 0.5, alpha=opacity)}
        c.setFillColor(color_map.get(color, color_map["gray"]))
        c.setFont("Helvetica-Bold", font_size)

        c.saveState()
        c.translate(pw / 2, ph / 2)
        c.rotate(45)
        c.drawCentredString(0, 0, text)
        c.restoreState()
        c.save()

        wm_reader = PdfReader(io.BytesIO(wm_buf.getvalue()))
        wm_page = wm_reader.pages[0]

        writer = PdfWriter()
        for page in reader.pages:
            page.merge_page(wm_page)
            writer.add_page(page)

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(file.filename).stem}_watermarked.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Watermark failed: {e}")
    finally:
        cleanup(src, wm_path)


# ─────────────────────────────────────────────────────────────────────────────
# 14. Fill & Sign  /api/pdf/fill-sign
# ─────────────────────────────────────────────────────────────────────────────
@router.post("/fill-sign", summary="Overlay a signature image onto a PDF page")
async def fill_sign(
    pdf_file: UploadFile = File(...),
    signature: UploadFile = File(...),
    page: int = Form(1),          # 1-indexed
    x: float = Form(100),         # points from left
    y: float = Form(100),         # points from bottom
    width: float = Form(150),
    height: float = Form(60),
):
    require_mime(pdf_file, "pdf")
    require_mime(signature, "image")

    src = await save_upload(pdf_file, ".pdf")
    sig_path = await save_upload(signature, ".png")

    try:
        reader = PdfReader(str(src))
        total = len(reader.pages)
        page_idx = max(0, min(page - 1, total - 1))
        target = reader.pages[page_idx]

        pw = float(target.mediabox.width)
        ph = float(target.mediabox.height)

        # Convert sig to PNG with alpha
        sig_img = PILImage.open(str(sig_path)).convert("RGBA")
        sig_buf = io.BytesIO()
        sig_img.save(sig_buf, format="PNG")

        # Build overlay PDF
        overlay_buf = io.BytesIO()
        c = rl_canvas.Canvas(overlay_buf, pagesize=(pw, ph))
        sig_buf.seek(0)
        c.drawImage(sig_buf, x, y, width=width, height=height, mask="auto")
        c.save()

        overlay_reader = PdfReader(io.BytesIO(overlay_buf.getvalue()))
        target.merge_page(overlay_reader.pages[0])

        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)

        buf = io.BytesIO()
        writer.write(buf)
        return Response(
            content=buf.getvalue(),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{Path(pdf_file.filename).stem}_signed.pdf"'},
        )
    except Exception as e:
        raise tool_error(f"Fill & sign failed: {e}")
    finally:
        cleanup(src, sig_path)
